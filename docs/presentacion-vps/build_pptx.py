"""
Presentación unificada para VPs — Plataforma de Gestión de Portfolio (PAE).
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Brandbook PAE
PAE_RED   = RGBColor(0xC8, 0x10, 0x2E)
PAE_BLUE  = RGBColor(0x00, 0x3D, 0xA5)
PAE_GREEN = RGBColor(0x00, 0x84, 0x3D)
BG_SOFT   = RGBColor(0xEB, 0xF0, 0xF7)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
TXT_1     = RGBColor(0x33, 0x33, 0x33)
TXT_2     = RGBColor(0x66, 0x66, 0x66)
TXT_3     = RGBColor(0x99, 0x99, 0x99)
BORDER    = RGBColor(0xDD, 0xDD, 0xDD)

FONT = "Inter"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG_SOFT
    return s


def add_text(slide, x, y, w, h, text, *, size=14, bold=False, color=TXT_1,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.05)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb


def add_rect(slide, x, y, w, h, fill=WHITE, line=None, line_w=0.75):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(line_w)
    sh.shadow.inherit = False
    return sh


def add_round(slide, x, y, w, h, fill=WHITE, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sh.adjustments[0] = 0.08
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
    return sh


def header(slide, title, subtitle=None):
    add_rect(slide, 0, 0, SW, Inches(0.12), fill=PAE_BLUE)
    add_text(slide, Inches(0.5), Inches(0.35), Inches(12.3), Inches(0.55),
             title, size=26, bold=True, color=PAE_BLUE)
    if subtitle:
        add_text(slide, Inches(0.5), Inches(0.95), Inches(12.3), Inches(0.4),
                 subtitle, size=13, color=TXT_2)
    add_rect(slide, Inches(0.5), Inches(1.38), Inches(12.3), Emu(12700), fill=BORDER)


def footer(slide, page_num):
    add_text(slide, Inches(0.5), Inches(7.12), Inches(6), Inches(0.3),
             "Pan American Energy  |  Plataforma de Gestión de Portfolio",
             size=9, color=TXT_3)
    add_text(slide, Inches(11.3), Inches(7.12), Inches(1.5), Inches(0.3),
             f"{page_num}", size=9, color=TXT_3, align=PP_ALIGN.RIGHT)


# ---------- SLIDE 1: COVER ----------
s = add_slide()
add_rect(s, 0, 0, SW, SH, fill=WHITE)
add_rect(s, 0, 0, Inches(0.25), SH, fill=PAE_BLUE)
add_rect(s, 0, SH - Inches(0.25), SW, Inches(0.25), fill=PAE_RED)

add_text(s, Inches(1.0), Inches(1.3), Inches(11), Inches(0.5),
         "PAN AMERICAN ENERGY", size=14, bold=True, color=PAE_RED)
add_text(s, Inches(1.0), Inches(2.0), Inches(11), Inches(1.3),
         "Plataforma de Gestión\nde Portfolio", size=48, bold=True, color=PAE_BLUE)
add_text(s, Inches(1.0), Inches(3.9), Inches(11), Inches(1.2),
         "Una herramienta centralizada para proponer, evaluar, aprobar y dar\n"
         "seguimiento a las iniciativas de transformación digital de PAE.",
         size=16, color=TXT_2)

add_rect(s, Inches(1.0), Inches(5.4), Inches(3.8), Emu(18000), fill=PAE_BLUE)
add_text(s, Inches(1.0), Inches(5.55), Inches(6), Inches(0.5),
         "Presentación a Vicepresidencias", size=16, bold=True, color=PAE_BLUE)
add_text(s, Inches(1.0), Inches(6.05), Inches(6), Inches(0.4),
         "Área de Transformación Digital  |  Abril 2026",
         size=12, color=TXT_2)


# ---------- SLIDE 2: PROBLEMA ----------
s = add_slide(); header(s, "El problema actual",
    "Por qué existe esta plataforma")
items = [
    ("Sin centralización",
     "Documentos Word sueltos, emails dispersos y reuniones sin registro estandarizado.\n"
     "Cada área propone y gestiona iniciativas con su propio formato.", PAE_RED),
    ("Sin trazabilidad",
     "Los formularios se pierden, se duplican o se desactualizan. No hay historial\n"
     "de cambios, ni versión final firmada, ni una fuente única de verdad.", PAE_RED),
    ("Sin visibilidad ejecutiva",
     "Los VPs y sponsors no tienen una vista consolidada del portfolio:\n"
     "estado, valor esperado, gasto acumulado, gates pendientes.", PAE_RED),
    ("Proceso de aprobación informal",
     "Las decisiones de seguir/pausar/rechazar una iniciativa se toman en reuniones\n"
     "sin registro estructurado del feedback de cada aprobador.", PAE_RED),
]
y = Inches(1.7)
for title, desc, c in items:
    add_rect(s, Inches(0.5), y, Inches(12.3), Inches(1.15), fill=WHITE, line=BORDER)
    add_rect(s, Inches(0.5), y, Emu(50000), Inches(1.15), fill=c)
    add_text(s, Inches(0.75), y + Inches(0.12), Inches(11.8), Inches(0.4),
             title, size=15, bold=True, color=c)
    add_text(s, Inches(0.75), y + Inches(0.48), Inches(11.8), Inches(0.65),
             desc, size=11, color=TXT_2)
    y += Inches(1.27)
footer(s, 2)


# ---------- SLIDE 3: OBJETIVO ----------
s = add_slide(); header(s, "Objetivo de la plataforma",
    "Digitalizar y estandarizar el ciclo de vida de toda iniciativa de transformación")
cells = [
    ("Digitalizar formularios", "5 formularios digitales con wizard, carry-over automático y validación por sección."),
    ("Gestionar gateways", "3 gateways formales con unanimidad, 5 estados posibles y feedback trazable."),
    ("Generar documentos", "PPTX, PDF, Excel y notas de prensa generados con un click desde la app."),
    ("Dashboards por rol", "Vistas diferenciadas para PO, Sponsor, VP, Business Owner y Área Transformación."),
    ("Trazabilidad total", "Historial de cambios, versiones firmadas (VF) y registro completo de aprobaciones."),
    ("Visión ejecutiva", "KPIs consolidados del portfolio: valor, ROI, gasto, gates pendientes, distribución por etapa."),
]
col_w = Inches(4.0); col_h = Inches(2.5); gap = Inches(0.15)
x0 = Inches(0.5); y0 = Inches(1.7)
for i, (t, d) in enumerate(cells):
    cx = x0 + (col_w + gap) * (i % 3)
    cy = y0 + (col_h + gap) * (i // 3)
    add_rect(s, cx, cy, col_w, col_h, fill=WHITE, line=BORDER)
    add_rect(s, cx, cy, col_w, Emu(45000), fill=PAE_BLUE)
    add_text(s, cx + Inches(0.2), cy + Inches(0.2), col_w - Inches(0.4), Inches(0.5),
             t, size=15, bold=True, color=PAE_BLUE)
    add_text(s, cx + Inches(0.2), cy + Inches(0.8), col_w - Inches(0.4), col_h - Inches(1.0),
             d, size=11, color=TXT_2)
footer(s, 3)


# ---------- SLIDE 4: USUARIOS Y ROLES ----------
s = add_slide(); header(s, "Usuarios y roles",
    "~200 empleados internos de PAE. Dos capas de permisos combinadas.")

# Capa 1
add_text(s, Inches(0.5), Inches(1.7), Inches(12), Inches(0.4),
         "Capa 1 — Perfil global", size=14, bold=True, color=PAE_BLUE)
profs = [
    ("User", "Empleado estándar. Propone iniciativas y participa según su rol por iniciativa."),
    ("Área Transformación", "Vista global del portfolio, administra permisos y facilita gateways."),
    ("Admin", "Gestiona usuarios, permisos y configuración global."),
    ("VP / Sponsor", "Detectado automáticamente. Vista ejecutiva y aprobaciones."),
]
y = Inches(2.15)
for t, d in profs:
    add_rect(s, Inches(0.5), y, Inches(12.3), Inches(0.55), fill=WHITE, line=BORDER)
    add_text(s, Inches(0.75), y + Inches(0.1), Inches(3), Inches(0.4),
             t, size=12, bold=True, color=PAE_BLUE)
    add_text(s, Inches(3.8), y + Inches(0.1), Inches(9), Inches(0.4),
             d, size=11, color=TXT_2)
    y += Inches(0.62)

# Capa 2
add_text(s, Inches(0.5), Inches(4.85), Inches(12), Inches(0.4),
         "Capa 2 — Rol por iniciativa", size=14, bold=True, color=PAE_BLUE)
roles = ["Promotor", "Líder de Dimensionamiento", "Product Owner", "Business Owner",
         "Sponsor", "Scrum Master", "Equipo"]
rx = Inches(0.5)
for r in roles:
    w = Inches(1.75)
    add_round(s, rx, Inches(5.35), w, Inches(0.5), fill=PAE_BLUE)
    add_text(s, rx, Inches(5.45), w, Inches(0.35),
             r, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rx += w + Inches(0.05)

add_text(s, Inches(0.5), Inches(6.1), Inches(12.3), Inches(0.9),
         "Row-level security: cada función de acceso a datos verifica permisos del usuario\n"
         "sobre cada iniciativa. Un usuario solo ve lo que le corresponde por su combinación de roles.",
         size=11, color=TXT_2)
footer(s, 4)


# ---------- SLIDE 5: CICLO DE VIDA ----------
s = add_slide(); header(s, "Ciclo de vida de una iniciativa",
    "3 etapas secuenciales con gateway + 2 ciclos anuales sin gateway")

stages = [
    ("Etapa 1", "Propuesta", "F1 + G1", PAE_BLUE),
    ("Etapa 2", "Dimensionamiento", "F2 + G2", PAE_BLUE),
    ("Etapa 3", "MVP", "F3 + G3", PAE_BLUE),
    ("Ciclo anual", "Visión Anual", "F4 (sin gateway)", PAE_GREEN),
    ("Ciclo anual", "Planificación", "F5 (sin gateway)", PAE_GREEN),
]
box_w = Inches(2.3); box_h = Inches(1.6); gap_x = Inches(0.15)
x = Inches(0.5); y = Inches(1.85)
for i, (tag, name, sub, c) in enumerate(stages):
    add_rect(s, x, y, box_w, box_h, fill=WHITE, line=c, line_w=1.5)
    add_rect(s, x, y, box_w, Inches(0.35), fill=c)
    add_text(s, x, y + Inches(0.05), box_w, Inches(0.3),
             tag, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, x, y + Inches(0.55), box_w, Inches(0.5),
             name, size=14, bold=True, color=c, align=PP_ALIGN.CENTER)
    add_text(s, x, y + Inches(1.1), box_w, Inches(0.4),
             sub, size=11, color=TXT_2, align=PP_ALIGN.CENTER)
    if i < len(stages) - 1:
        arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                   x + box_w + Emu(5000), y + Inches(0.6),
                                   gap_x - Emu(10000), Inches(0.4))
        arrow.fill.solid(); arrow.fill.fore_color.rgb = TXT_3
        arrow.line.fill.background()
    x += box_w + gap_x

# Camino A / B
add_text(s, Inches(0.5), Inches(3.8), Inches(12.3), Inches(0.4),
         "Dos caminos de entrada", size=14, bold=True, color=PAE_BLUE)
add_rect(s, Inches(0.5), Inches(4.25), Inches(6.0), Inches(1.2), fill=WHITE, line=BORDER)
add_text(s, Inches(0.7), Inches(4.35), Inches(5.6), Inches(0.4),
         "Camino A — Iniciativa nueva", size=12, bold=True, color=PAE_BLUE)
add_text(s, Inches(0.7), Inches(4.7), Inches(5.6), Inches(0.7),
         "F1 → G1 → F1 VF → F2 → G2 → F2 VF → F3 → G3 → F3 VF → F4 → F5",
         size=10, color=TXT_2)

add_rect(s, Inches(6.8), Inches(4.25), Inches(6.0), Inches(1.2), fill=WHITE, line=BORDER)
add_text(s, Inches(7.0), Inches(4.35), Inches(5.6), Inches(0.4),
         "Camino B — Iniciativa existente", size=12, bold=True, color=PAE_GREEN)
add_text(s, Inches(7.0), Inches(4.7), Inches(5.6), Inches(0.7),
         "Importar → F4 (ciclo anual) → F5 (ciclo anual)\n"
         "Muchas iniciativas no pasan por Etapa 1/2/3.", size=10, color=TXT_2)

# Estados de gateway
add_text(s, Inches(0.5), Inches(5.7), Inches(12.3), Inches(0.4),
         "5 estados posibles en cada gateway (unanimidad requerida)",
         size=13, bold=True, color=PAE_BLUE)
states = [("Aprobado", PAE_GREEN), ("Feedback", PAE_BLUE), ("Pausa", TXT_2),
          ("Rechazado", PAE_RED), ("Cambio de área", PAE_BLUE)]
sx = Inches(0.5)
for t, c in states:
    w = Inches(2.4)
    add_round(s, sx, Inches(6.2), w, Inches(0.55), fill=c)
    add_text(s, sx, Inches(6.3), w, Inches(0.35),
             t, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    sx += w + Inches(0.1)
footer(s, 5)


# ---------- SLIDE 6: LOS 5 FORMULARIOS ----------
s = add_slide(); header(s, "Los 5 formularios",
    "Cada formulario digital replica exactamente los Word originales de PAE")

forms = [
    ("F1 — Propuesta", "Gateway 1",
     "Problema, oportunidad, valor estratégico, alineación con objetivos PAE.\n"
     "9 secciones. Al aprobarse se firma como F1 VF y habilita F2.", PAE_BLUE),
    ("F2 — Dimensionamiento", "Gateway 2",
     "Factibilidad técnica, costos, equipo, riesgos, cronograma.\n"
     "Hereda automáticamente de F1 VF. Al aprobarse habilita F3.", PAE_BLUE),
    ("F3 — MVP", "Gateway 3",
     "Planificación y documentación del MVP. Resultados obtenidos vs esperados.\n"
     "Hereda de F2 VF. Al aprobarse pasa a ciclo anual F4/F5.", PAE_BLUE),
    ("F4 — Visión Anual", "Sin gateway",
     "Visión del producto para el año siguiente. Trabajo constante con seguimiento.\n"
     "Sponsor marca 'reviewed'. Se repite cada año.", PAE_GREEN),
    ("F5 — Planificación Anual", "Sin gateway",
     "Entregables concretos + indicadores de seguimiento del año.\n"
     "Hereda de F4 del mismo ciclo. Actualizaciones periódicas.", PAE_GREEN),
]
y = Inches(1.7)
for name, gate, desc, c in forms:
    add_rect(s, Inches(0.5), y, Inches(12.3), Inches(1.0), fill=WHITE, line=BORDER)
    add_rect(s, Inches(0.5), y, Emu(60000), Inches(1.0), fill=c)
    add_text(s, Inches(0.8), y + Inches(0.12), Inches(4), Inches(0.4),
             name, size=13, bold=True, color=c)
    add_round(s, Inches(4.3), y + Inches(0.15), Inches(1.4), Inches(0.35), fill=c)
    add_text(s, Inches(4.3), y + Inches(0.2), Inches(1.4), Inches(0.25),
             gate, size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.8), y + Inches(0.5), Inches(11.7), Inches(0.5),
             desc, size=10, color=TXT_2)
    y += Inches(1.07)
footer(s, 6)


# ---------- SLIDE 7: MOTOR DE FORMULARIOS ----------
s = add_slide(); header(s, "Motor de formularios",
    "Cómo se completan las iniciativas en la plataforma")

features = [
    ("Wizard por secciones",
     "Navegación vertical por steps. Cada sección se completa y guarda progresivamente.\n"
     "Auto-save cada 30s + al perder foco (debounce 1s)."),
    ("Carry-over automático",
     "F1 VF alimenta F2, F2 VF alimenta F3, F3 VF alimenta F4, F4 alimenta F5.\n"
     "Campos heredados aparecen en gris, editables al hacer click."),
    ("Subida de Word/Excel existentes",
     "La plataforma mapea automáticamente un documento existente contra el formulario\n"
     "digital y señala los campos faltantes."),
    ("Control de cambios",
     "Historial tipo Google Docs con autor y fecha. La versión inicial y la VF se\n"
     "conservan como versiones firmadas e inmutables."),
    ("Validación estricta",
     "TypeScript + Zod validan cada input. No se puede enviar un gateway con campos\n"
     "obligatorios incompletos."),
]
y = Inches(1.7)
for t, d in features:
    add_rect(s, Inches(0.5), y, Inches(12.3), Inches(0.95), fill=WHITE, line=BORDER)
    add_text(s, Inches(0.8), y + Inches(0.12), Inches(11.7), Inches(0.4),
             t, size=13, bold=True, color=PAE_BLUE)
    add_text(s, Inches(0.8), y + Inches(0.45), Inches(11.7), Inches(0.5),
             d, size=10, color=TXT_2)
    y += Inches(1.02)
footer(s, 7)


# ---------- SLIDE 8: GATEWAYS ----------
s = add_slide(); header(s, "Gateways de aprobación",
    "El mecanismo formal de decisión del portfolio")

add_rect(s, Inches(0.5), Inches(1.7), Inches(6.0), Inches(5.1), fill=WHITE, line=BORDER)
add_text(s, Inches(0.8), Inches(1.85), Inches(5.4), Inches(0.5),
         "Cómo funciona", size=15, bold=True, color=PAE_BLUE)
bullets = [
    "Cada aprobador registra su voto + comentarios.",
    "Se requiere unanimidad para aprobar.",
    "Feedback se consolida en un documento único.",
    "Prioridad de estados: reject > pause > area_change > feedback > pending > approved.",
    "Cada gateway define quién lidera la siguiente etapa y si se suman nuevos participantes.",
    "La versión firmada (VF) queda como registro inmutable.",
]
y = Inches(2.45)
for b in bullets:
    add_text(s, Inches(0.8), y, Inches(0.2), Inches(0.35), "•", size=14, bold=True, color=PAE_RED)
    add_text(s, Inches(1.0), y, Inches(5.3), Inches(0.7), b, size=11, color=TXT_1)
    y += Inches(0.65)

# 5 estados
add_rect(s, Inches(6.8), Inches(1.7), Inches(6.0), Inches(5.1), fill=WHITE, line=BORDER)
add_text(s, Inches(7.1), Inches(1.85), Inches(5.4), Inches(0.5),
         "5 estados posibles", size=15, bold=True, color=PAE_BLUE)
states_d = [
    ("Aprobado", "Todos los aprobadores votan a favor. Avanza a la siguiente etapa.", PAE_GREEN),
    ("Feedback", "Se solicitan ajustes. Vuelve al promotor y se re-envía.", PAE_BLUE),
    ("Pausa", "Se congela temporalmente. Puede retomarse más adelante.", TXT_2),
    ("Rechazado", "No avanza. Se cierra el caso con motivo.", PAE_RED),
    ("Cambio de área", "Reasigna la iniciativa a otra VP/área responsable.", PAE_BLUE),
]
y = Inches(2.45)
for t, d, c in states_d:
    add_round(s, Inches(7.1), y, Inches(1.5), Inches(0.4), fill=c)
    add_text(s, Inches(7.1), y + Inches(0.05), Inches(1.5), Inches(0.3),
             t, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(8.75), y + Inches(0.05), Inches(4.0), Inches(0.4),
             d, size=10, color=TXT_2)
    y += Inches(0.65)
footer(s, 8)


# ---------- SLIDE 9: GENERACIÓN DE DOCUMENTOS ----------
s = add_slide(); header(s, "Generación automática de documentos",
    "Desde cualquier formulario, con un click")

docs = [
    ("PPTX", "~10 slides\nbrandbook PAE", "pptxgenjs"),
    ("PDF", "Formulario completo\ncon formato institucional", "html2pdf.js"),
    ("XLSX", "Un tab por sección\ndel formulario", "SheetJS"),
    ("DOCX", "Notas de prensa\n(Working Backwards)\n+ minutas de gateway", "docx"),
]
x = Inches(0.5); y = Inches(1.9)
w = Inches(3.0); h = Inches(2.7); gap = Inches(0.1)
for t, d, lib in docs:
    add_rect(s, x, y, w, h, fill=WHITE, line=BORDER)
    add_rect(s, x, y, w, Inches(0.8), fill=PAE_BLUE)
    add_text(s, x, y + Inches(0.15), w, Inches(0.55),
             t, size=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, x, y + Inches(1.05), w, Inches(1.0),
             d, size=12, color=TXT_1, align=PP_ALIGN.CENTER)
    add_text(s, x, y + Inches(2.15), w, Inches(0.35),
             lib, size=10, color=TXT_3, align=PP_ALIGN.CENTER)
    x += w + gap

add_rect(s, Inches(0.5), Inches(5.0), Inches(12.3), Inches(1.8), fill=WHITE, line=BORDER)
add_text(s, Inches(0.8), Inches(5.15), Inches(11.7), Inches(0.4),
         "Ventaja clave: generación 100% en el browser", size=14, bold=True, color=PAE_GREEN)
add_text(s, Inches(0.8), Inches(5.55), Inches(11.7), Inches(1.2),
         "Todos los documentos se generan del lado del cliente, sin necesidad de servidor de\n"
         "generación de archivos. En MVP se descargan localmente; en fase backend (post-MVP)\n"
         "se archivarán automáticamente en SharePoint dentro de la carpeta de la iniciativa.",
         size=11, color=TXT_2)
footer(s, 9)


# ---------- SLIDE 10: DASHBOARDS POR ROL ----------
s = add_slide(); header(s, "Dashboards por rol",
    "Cada usuario ve lo que necesita para su rol, sin ruido")

dash = [
    ("PO / Líder / Promotor", PAE_BLUE,
     ["Iniciativas propias y participadas",
      "Acciones pendientes + timeline",
      "Edición de formularios",
      "Estado actual de cada gateway"]),
    ("Business Owner / Sponsor", PAE_BLUE,
     ["Aprobaciones en gateways",
      "Seguimiento de iniciativas patrocinadas",
      "Proponer cambios / feedback",
      "Registro de decisiones"]),
    ("VP", PAE_RED,
     ["Portfolio completo con filtros",
      "KPIs: valor, ROI, gasto, gates",
      "Ranking de iniciativas por ROI",
      "Próximos hitos y gates pendientes"]),
    ("Área Transformación", PAE_GREEN,
     ["Overview de todas las VPs",
      "Administración de usuarios y roles",
      "Métricas globales del portfolio",
      "Facilitación de gateways"]),
]
x = Inches(0.5); y = Inches(1.9)
w = Inches(3.0); h = Inches(4.7); gap = Inches(0.1)
for title, c, bs in dash:
    add_rect(s, x, y, w, h, fill=WHITE, line=BORDER)
    add_rect(s, x, y, w, Inches(0.7), fill=c)
    add_text(s, x + Inches(0.15), y + Inches(0.15), w - Inches(0.3), Inches(0.4),
             title, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    by = y + Inches(0.95)
    for b in bs:
        add_text(s, x + Inches(0.2), by, Inches(0.15), Inches(0.3),
                 "•", size=12, bold=True, color=c)
        add_text(s, x + Inches(0.4), by, w - Inches(0.5), Inches(0.8),
                 b, size=10, color=TXT_1)
        by += Inches(0.75)
    x += w + gap
footer(s, 10)


# ---------- SLIDE 11: MÉTRICAS CONSOLIDADAS ----------
s = add_slide(); header(s, "Vista ejecutiva del portfolio",
    "Lo que un VP ve cuando entra a la plataforma")

kpis = [("Iniciativas activas", "#"), ("Valor esperado total", "USD"),
        ("Gasto acumulado", "USD"), ("Gates pendientes", "#"),
        ("ROI promedio", "x"), ("En progreso / Pausa / Rechazo", "#/#/#")]
x = Inches(0.5); y = Inches(1.85)
w = Inches(4.0); h = Inches(1.3); gx = Inches(0.15); gy = Inches(0.15)
for i, (t, unit) in enumerate(kpis):
    cx = x + (w + gx) * (i % 3)
    cy = y + (h + gy) * (i // 3)
    add_rect(s, cx, cy, w, h, fill=WHITE, line=BORDER)
    add_rect(s, cx, cy, Emu(50000), h, fill=PAE_BLUE)
    add_text(s, cx + Inches(0.25), cy + Inches(0.15), w - Inches(0.3), Inches(0.4),
             t, size=12, bold=True, color=TXT_1)
    add_text(s, cx + Inches(0.25), cy + Inches(0.55), w - Inches(0.3), Inches(0.6),
             unit, size=26, bold=True, color=PAE_BLUE)

add_rect(s, Inches(0.5), Inches(4.95), Inches(12.3), Inches(1.9), fill=WHITE, line=BORDER)
add_text(s, Inches(0.8), Inches(5.1), Inches(11.7), Inches(0.4),
         "Otros elementos del dashboard", size=14, bold=True, color=PAE_BLUE)
items = [
    "Distribución por etapa (F1 / F2 / F3 / F4 / F5) con cantidad de iniciativas en cada una.",
    "Corrientes de valor — beneficio bruto proyectado del año consolidado.",
    "Ranking de iniciativas por ROI descendente.",
    "Próximos eventos: gateways programados, sprint reviews e hitos.",
]
ry = Inches(5.5)
for it in items:
    add_text(s, Inches(0.85), ry, Inches(0.2), Inches(0.3), "•", size=12, bold=True, color=PAE_RED)
    add_text(s, Inches(1.05), ry, Inches(11.5), Inches(0.35), it, size=10, color=TXT_2)
    ry += Inches(0.33)
footer(s, 11)


# ---------- SLIDE 12: MVP (Fase actual) ----------
s = add_slide(); header(s, "Alcance del MVP — Fase actual",
    "Lo que la plataforma hace HOY (frontend-first, fases 2–4)")

mvp = [
    "5 formularios digitales con wizard por secciones (F1, F2, F3, F4, F5)",
    "Carry-over automático entre formularios + campos heredados editables",
    "3 gateways con unanimidad, 5 estados posibles y feedback consolidado",
    "Versión Firmada (VF) inmutable al aprobarse cada gateway",
    "Generación automática de PPTX, PDF, XLSX y DOCX (browser-side)",
    "Dashboards diferenciados por rol (PO, Sponsor, VP, Área Transformación)",
    "Vista ejecutiva del portfolio con KPIs, ROI, ranking y próximos eventos",
    "Row-level security: cada usuario solo ve lo que le corresponde",
    "Auto-save cada 30s + al perder foco — cero pérdida de trabajo",
    "Control de cambios con historial (autor + fecha) tipo Google Docs",
    "Validación estricta con TypeScript + Zod en cada input",
    "Notificaciones in-app con datos persistidos",
    "Dos caminos: iniciativa nueva (F1→G1→...) o existente (importar → F4/F5)",
    "Persistencia en localStorage (replica el modelo de datos final)",
    "Mock auth con selector de usuario para demo y testing interno",
]
# two columns
col_w = Inches(6.0); gap = Inches(0.3)
x1, x2 = Inches(0.5), Inches(0.5) + col_w + gap
y0 = Inches(1.85)
half = (len(mvp) + 1) // 2
for i, item in enumerate(mvp):
    cx = x1 if i < half else x2
    cy = y0 + Inches(0.32) * (i % half)
    add_rect(s, cx, cy, Emu(50000), Inches(0.25), fill=PAE_GREEN)
    add_text(s, cx + Inches(0.15), cy - Inches(0.03), col_w - Inches(0.2), Inches(0.35),
             item, size=10, color=TXT_1)
footer(s, 12)


# ---------- SLIDE 13: PUNTOS DE MEJORA — POST-MVP + BACKEND ----------
s = add_slide(); header(s, "Puntos de mejora — Post-MVP + Backend",
    "Todo lo que queda fuera del MVP actual y requiere la fase backend (Fase 5)")

add_text(s, Inches(0.5), Inches(1.7), Inches(12.3), Inches(0.4),
         "Fase 5 — Backend e infraestructura", size=14, bold=True, color=PAE_RED)
backend = [
    ("Base de datos PostgreSQL + Prisma",
     "Reemplaza la persistencia localStorage. Habilita concurrencia real entre ~200 usuarios."),
    ("SSO con Microsoft Entra ID",
     "Reemplaza el mock auth. Login con cuenta corporativa PAE y permisos federados."),
    ("Integración SharePoint",
     "Carpetas automáticas por iniciativa. Archivo de PPTX/PDF/XLSX/DOCX y de docs adicionales."),
    ("Hosting AWS + red interna PAE",
     "Deploy en infraestructura PAE. Acceso solo por VPN / red interna, nada expuesto a internet."),
    ("Notificaciones por Outlook",
     "Complementa las in-app con emails automáticos en eventos clave del ciclo."),
]
y = Inches(2.1)
for t, d in backend:
    add_rect(s, Inches(0.5), y, Inches(12.3), Inches(0.7), fill=WHITE, line=BORDER)
    add_rect(s, Inches(0.5), y, Emu(50000), Inches(0.7), fill=PAE_RED)
    add_text(s, Inches(0.8), y + Inches(0.08), Inches(5.2), Inches(0.3),
             t, size=11, bold=True, color=PAE_RED)
    add_text(s, Inches(6.1), y + Inches(0.08), Inches(6.6), Inches(0.6),
             d, size=10, color=TXT_2)
    y += Inches(0.77)

# Post-MVP funcional
add_text(s, Inches(0.5), Inches(5.95), Inches(12.3), Inches(0.4),
         "Funcionalidad post-MVP", size=14, bold=True, color=PAE_BLUE)
post = [
    "Vision House (objetivos estratégicos de PAE vinculados a iniciativas)",
    "Cash outflow automatizado desde SAP",
    "Métricas de performance integradas con Jira",
    "Dashboards avanzados (proyecciones, cohortes, drill-down por VP)",
    "Integraciones con otros sistemas internos (Power BI, Teams)",
]
y = Inches(6.35)
for p in post:
    add_text(s, Inches(0.85), y, Inches(0.2), Inches(0.3), "•", size=12, bold=True, color=PAE_BLUE)
    add_text(s, Inches(1.05), y, Inches(11.5), Inches(0.35), p, size=10, color=TXT_2)
    y += Inches(0.22)
footer(s, 13)


# ---------- SLIDE 14: STACK TÉCNICO ----------
s = add_slide(); header(s, "Stack técnico",
    "Fase actual vs Fase 5 (backend)")

add_rect(s, Inches(0.5), Inches(1.7), Inches(6.1), Inches(5.1), fill=WHITE, line=BORDER)
add_rect(s, Inches(0.5), Inches(1.7), Inches(6.1), Inches(0.5), fill=PAE_GREEN)
add_text(s, Inches(0.5), Inches(1.8), Inches(6.1), Inches(0.35),
         "Fase actual (2–4) — Frontend-first", size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
rows1 = [
    ("Frontend", "Next.js 15+ App Router, TypeScript estricto"),
    ("Estilos", "Tailwind CSS con theme PAE (brandbook)"),
    ("Persistencia", "localStorage (replica el modelo final)"),
    ("Auth", "Mock: selector de usuario sin password"),
    ("Documentos", "pptxgenjs, SheetJS, html2pdf.js, docx"),
    ("Validación", "Zod en todos los inputs"),
    ("Notificaciones", "In-app con datos en localStorage"),
]
y = Inches(2.35)
for k, v in rows1:
    add_text(s, Inches(0.8), y, Inches(2.0), Inches(0.35), k, size=11, bold=True, color=PAE_GREEN)
    add_text(s, Inches(2.8), y, Inches(3.7), Inches(0.6), v, size=10, color=TXT_1)
    y += Inches(0.6)

add_rect(s, Inches(6.75), Inches(1.7), Inches(6.1), Inches(5.1), fill=WHITE, line=BORDER)
add_rect(s, Inches(6.75), Inches(1.7), Inches(6.1), Inches(0.5), fill=PAE_RED)
add_text(s, Inches(6.75), Inches(1.8), Inches(6.1), Inches(0.35),
         "Fase 5 — Backend e infraestructura", size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
rows2 = [
    ("Frontend", "(se conserva sin cambios mayores)"),
    ("Base de datos", "PostgreSQL + Prisma ORM"),
    ("Auth", "SSO con Microsoft Entra ID"),
    ("Storage docs", "SharePoint (carpetas automáticas)"),
    ("Hosting", "AWS sobre red interna PAE / VPN"),
    ("Notificaciones", "Outlook + in-app"),
    ("Integraciones", "SAP, Jira, Power BI, Teams"),
]
y = Inches(2.35)
for k, v in rows2:
    add_text(s, Inches(7.05), y, Inches(2.0), Inches(0.35), k, size=11, bold=True, color=PAE_RED)
    add_text(s, Inches(9.05), y, Inches(3.7), Inches(0.6), v, size=10, color=TXT_1)
    y += Inches(0.6)
footer(s, 14)


# ---------- SLIDE 15: RESULTADO ESPERADO ----------
s = add_slide(); header(s, "Resultado esperado",
    "Qué gana PAE con la plataforma en producción")

kpis = [
    ("100%", "Trazabilidad", "Ciclo completo desde la idea hasta F5 con historial de cambios y VFs firmadas.", PAE_GREEN),
    ("−70%", "Tiempo de gestión", "Eliminando documentos duplicados, emails, re-trabajo manual y re-ingreso de datos.", PAE_BLUE),
    ("Real-time", "Visibilidad ejecutiva", "Estado del portfolio actualizado para VPs y sponsors en todo momento.", PAE_BLUE),
    ("200+", "Usuarios", "Plataforma preparada para escalar a toda la organización PAE.", PAE_RED),
]
x = Inches(0.5); y = Inches(1.9); w = Inches(3.0); h = Inches(4.7); gap = Inches(0.1)
for big, t, d, c in kpis:
    add_rect(s, x, y, w, h, fill=WHITE, line=BORDER)
    add_rect(s, x, y, w, Inches(1.8), fill=c)
    add_text(s, x, y + Inches(0.35), w, Inches(1.1),
             big, size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, x, y + Inches(2.0), w, Inches(0.5),
             t, size=15, bold=True, color=c, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.25), y + Inches(2.6), w - Inches(0.5), Inches(1.9),
             d, size=11, color=TXT_2, align=PP_ALIGN.CENTER)
    x += w + gap
footer(s, 15)


# ---------- SLIDE 16: PRÓXIMOS PASOS ----------
s = add_slide(); header(s, "Próximos pasos",
    "De dónde venimos y qué sigue")

steps = [
    ("01", "Validar esta visión con las VPs", "Reunión actual. Recoger feedback sobre alcance y prioridades."),
    ("02", "Cerrar el MVP frontend", "Completar el wizard de los 5 formularios y refinar dashboards por rol."),
    ("03", "Piloto interno con Área Transformación", "Probar con iniciativas reales en modo localStorage."),
    ("04", "Alineación con IT", "Confirmar accesos para Entra ID (SSO), SharePoint API y hosting AWS interno."),
    ("05", "Iniciar Fase 5 — Backend", "Migración a PostgreSQL + Prisma, integración SharePoint, deploy en red PAE."),
    ("06", "Rollout escalonado a toda la organización", "Onboarding por VP con soporte del Área Transformación."),
]
y = Inches(1.85)
for num, t, d in steps:
    add_rect(s, Inches(0.5), y, Inches(12.3), Inches(0.78), fill=WHITE, line=BORDER)
    add_rect(s, Inches(0.5), y, Inches(0.9), Inches(0.78), fill=PAE_BLUE)
    add_text(s, Inches(0.5), y + Inches(0.22), Inches(0.9), Inches(0.4),
             num, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(1.6), y + Inches(0.1), Inches(11), Inches(0.35),
             t, size=13, bold=True, color=PAE_BLUE)
    add_text(s, Inches(1.6), y + Inches(0.42), Inches(11), Inches(0.35),
             d, size=10, color=TXT_2)
    y += Inches(0.85)
footer(s, 16)


# ---------- SLIDE 17: CIERRE ----------
s = add_slide()
add_rect(s, 0, 0, SW, SH, fill=WHITE)
add_rect(s, 0, 0, SW, Inches(0.25), fill=PAE_BLUE)
add_rect(s, 0, SH - Inches(0.25), SW, Inches(0.25), fill=PAE_RED)

add_text(s, Inches(1), Inches(2.0), Inches(11), Inches(0.5),
         "PAN AMERICAN ENERGY", size=14, bold=True, color=PAE_RED, align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(2.7), Inches(11), Inches(1.2),
         "Plataforma de Gestión de Portfolio", size=40, bold=True, color=PAE_BLUE, align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(4.2), Inches(11), Inches(0.6),
         "Gracias", size=28, bold=True, color=TXT_1, align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(5.0), Inches(11), Inches(0.5),
         "Preguntas y comentarios", size=16, color=TXT_2, align=PP_ALIGN.CENTER)
add_text(s, Inches(1), Inches(6.3), Inches(11), Inches(0.4),
         "Área de Transformación Digital  |  Abril 2026",
         size=11, color=TXT_3, align=PP_ALIGN.CENTER)


out = "/Users/martupaz/iniciativas/docs/presentacion-vps/Plataforma_Gestion_Portfolio_VPs.pptx"
prs.save(out)
print("OK:", out)
